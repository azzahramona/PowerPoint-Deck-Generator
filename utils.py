# utils.py
import time
import json
import numpy as np
import pandas as pd
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

## ppt format function
def format_title(slide, text, alignment, font_name, font_size, font_bold = False, font_color = RGBColor(0, 0, 0),left=Pt(75), top=Pt(25), width=Pt(850), height=Pt(70)):
    title_shape = slide.shapes.add_textbox(left=left, top=top, width=width, height=height)
    title_text_frame = title_shape.text_frame
    title_text_frame.text = text
    title_text_frame.word_wrap = True
    for paragraph in title_text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = font_name   
            run.font.bold = font_bold
            run.font.color.rgb = font_color
            run.font.size = Pt(font_size)
    return title_shape

def format_shape(slide, shape, left, top, width, height, color = RGBColor(255,255,255), outline_color = RGBColor(0,0,0), outline_width = Inches(0.005)):
    # Add the shape to the slide
    shape = slide.shapes.add_shape(shape, left, top, width, height)  # 9 represents the MSO_SHAPE.OVAL type
    
    # (Optional) Customize the circle's appearance
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color  # Set fill color to blue (RGB)

    line = shape.line
    line.width = outline_width  # Set line width
    line.color.rgb = outline_color  # Set line color to red (RGB)

    return shape

# Function to create bar clustered chart
def bar_clustered_chart(slide, df, x, y, cx, cy, legend=True, data_show=False, value_col=''):
    chart_data = CategoryChartData()
    for i in df.index:
        chart_data.add_category(i)
    for col in df.columns:
        chart_data.add_series(col, df[col].values)

    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart

    if legend:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.font.size = Pt(12)
    else:
        chart.has_legend = False

    if data_show:
        for series in chart.series:
            for i, val in enumerate(series.values):
                if val == 0:
                    series.points[i].data_label.has_text_frame = True
                series.data_labels.show_value = True
                series.data_labels.font.size = Pt(10)
                series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    chart.has_title = False

    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.axis_title.text_frame.text = str(value_col).title()
    chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(12)

    value_axis = chart.value_axis
    category_axis = chart.category_axis
    value_axis.has_major_gridlines = False
    value_axis.has_minor_gridlines = False
    category_axis.has_major_gridlines = False
    category_axis.has_minor_gridlines = False

    category_axis.tick_labels.font.size = Pt(10)

    return chart
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Pt
from pptx.dml.color import RGBColor
def bubble_chart(
    slide, df, x, y, cx, cy,
    x_col='x', y_col='y',
    legend=True, data_show=True
):
    """
    Create a bubble chart on the given slide.

    Parameters:
    - slide: pptx Slide object
    - df: DataFrame with columns ['Demographics', x_col, y_col, 'Size', 'colors']
          'colors' should be list/tuple RGB like [R,G,B]
    - x, y, cx, cy: float - position and size (use Inches or Pt)
    - x_col: str - column name for X values
    - y_col: str - column name for Y values
    - legend: bool - show legend
    - data_show: bool - show data labels (Demographics names)
    """

    # Prepare chart data
    chart_data = BubbleChartData()
    for i, row in df.iterrows():
        series = chart_data.add_series(row['Demographics'])
        series.add_data_point(row[x_col], row[y_col], row['Size'])

    # Add chart to slide
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
    ).chart

    # Set bubble colors
    for i, row in df.iterrows():
        fill = chart.series[i].points[0].format.fill
        fill.solid()
        color = RGBColor(*row['colors'])  # Expecting [R, G, B]
        fill.fore_color.rgb = color

    # Legend settings
    if legend:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.font.size = Pt(10)
    else:
        chart.has_legend = False

    # Data labels with Demographics names
    if data_show:
        for series in chart.series:
            point = series.points[0]
            point.data_label.text_frame.text = series.name
            for run in point.data_label.text_frame.paragraphs[0].runs:
                run.font.size = Pt(10)
            point.data_label.position = XL_LABEL_POSITION.ABOVE

    chart.has_title = False

    # Remove gridlines
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False

    # Axis titles and label sizes
    chart.value_axis.axis_title.text_frame.text = y_col.title()
    chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.value_axis.tick_labels.font.size = Pt(10)

    chart.category_axis.axis_title.text_frame.text = x_col.title()
    chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.category_axis.tick_labels.font.size = Pt(10)

    return chart
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_MARKER_STYLE
from pptx.dml.color import RGBColor
from pptx.util import Pt

def scatter_plot(
    slide, df, x_col, y_col, label_col, x, y, cx, cy,
    point_size=10, point_color=RGBColor(0, 0, 0),
    fontsize=Pt(9), legend=False, data_label=True
):
    """
    Create a scatter plot chart on the slide.

    Parameters:
    - slide: pptx Slide object
    - df: DataFrame with data
    - x_col: str - column name for X-axis values
    - y_col: str - column name for Y-axis values
    - label_col: str - column name for series labels
    - x, y, cx, cy: float - position and size of chart
    - point_size: int - size of the scatter points
    - point_color: RGBColor - color of the scatter points
    - fontsize: Pt - font size for labels and legend
    - legend: bool - show legend or not
    - data_label: bool - show data labels on points
    """

    # Prepare chart data
    chart_data = XyChartData()
    # Change this inside your scatter_plot function:
    for i in range(len(df)):
        label = df.loc[i, label_col] if label_col in df.columns else f"Point {i+1}"
        series_i = chart_data.add_series(label)
        series_i.add_data_point(df.loc[i, x_col], df.loc[i, y_col])


    # Add chart to slide
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
    ).chart

    # Remove Gridlines
    value_axis = chart.value_axis
    category_axis = chart.category_axis
    value_axis.has_major_gridlines = False
    value_axis.has_minor_gridlines = False
    category_axis.has_major_gridlines = False
    category_axis.has_minor_gridlines = False

    # Legend settings
    if legend:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = fontsize
    else:
        chart.has_legend = False

    # Data labels on points
    if data_label:
        for series in chart.series:
            point = series.points[0]
            point.data_label.text_frame.text = series.name
            for run in point.data_label.text_frame.paragraphs[0].runs:
                run.font.size = Pt(10)
            point.data_label.position = XL_LABEL_POSITION.ABOVE

    # Axis titles and formatting
    value_axis.axis_title.text_frame.text = y_col.title()
    value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(12)
    value_axis.tick_labels.font.size = fontsize

    category_axis.axis_title.text_frame.text = x_col.title()
    category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(12)
    category_axis.tick_labels.font.size = fontsize

    # Customize marker style and color
    for series in chart.series:
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.size = point_size
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = point_color
        series.marker.format.line.fill.background()  # make outline transparent

    chart.has_title = False

    return chart

def column_stacked_chart(slide, df, x, y, cx, cy):
    from pptx.chart.data import CategoryChartData

    # Transpose: Brands become categories, months become series
    df = df.copy()

    chart_data = CategoryChartData()
    chart_data.categories = list(df.index)  # Brand names on X-axis

    for col in df.columns:
        # Each column (month) becomes a stacked series
        chart_data.add_series(str(col), df[col].values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
    ).chart

    # Formatting
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.visible = False

    for series in chart.series:
        series.data_labels.show_value = True
        series.data_labels.font.size = Pt(9)
        series.data_labels.number_format = '0%'
        series.data_labels.position = XL_LABEL_POSITION.INSIDE_END

    return chart
def bar_stacked_chart(slide,df,x,y,cx,cy,percentage = True, data_show = True, legend = True, axis = True, fontsize = Pt(12), fontsize_label = Pt(10)):
    df.fillna(0, inplace = True) #fill nan
    # Define chart data
    chart_data = CategoryChartData()
    for i in df.columns:
        chart_data.add_category(i)
    for j, row in df.iterrows():
        chart_data.add_series(j, row.values)
    
    if percentage:
        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data).chart
    else:
        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data).chart
    
    if legend:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = fontsize
    
    chart.plots[0].gap_width = 50

    category_axis = chart.category_axis
    if axis:
        category_axis.has_major_gridlines = False
        category_axis.has_minor_gridlines = False
        category_axis.tick_labels.font.size = fontsize
    else:
        category_axis.visible = False
        
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = False
    value_axis.has_minor_gridlines = False
    value_axis.visible = False

    if data_show:
        for series in chart.plots[0].series:
            for i, val in enumerate(series.values):
                if val == 0:
                    series.points[i].data_label.has_text_frame = True
                series.data_labels.show_value = True
                series.data_labels.font.size = fontsize_label
                series.data_labels.number_format = '0%' if percentage else '#,##0'
                series.data_labels.position = XL_LABEL_POSITION.INSIDE_BASE
    return chart
def line_marker_chart(slide,df,x,y,cx,cy, legend = True, legend_position = XL_LEGEND_POSITION.RIGHT, data_show = False, chart_title = False, title ="", fontsize = Pt(9), fontsize_title = Pt(14), percentage = False, line_width = Pt(1)):
    df.fillna(0, inplace = True)
    # Define chart data
    chart_data = CategoryChartData()
    # for i in df.columns:
    #     chart_data.add_category(i)
    # for j, row in df.iterrows():
    #     chart_data.add_series(j, row.values)
    for i in df.columns:
        chart_data.add_category(i)
    for j, row in df.iterrows():
        chart_data.add_series(j,np.where(row.values == 0, None, row.values))
    
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data).chart
    
    # Add legend
    if legend:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = legend_position
        chart.legend.font.size = fontsize  # assuming Pt is imported from pptx.util
    else:
        chart.has_legend = False
    
    if data_show:
        for series in chart.plots[0].series:
            for i, val in enumerate(series.values):
                if val == 0:
                    series.points[i].data_label.has_text_frame = True
                series.data_labels.show_value = True
                series.data_labels.font.size = fontsize
                # series.data_labels.number_format = '0.00%'
                series.data_labels.position = XL_LABEL_POSITION.ABOVE
                
    # Change line point to all circle
    # Iterate through each series and set marker style to circle
    for series in chart.series:
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.format.line.width = line_width # custom line width

    # Customize y-axis format
    chart.value_axis.tick_labels.font.size = fontsize  # Set font size for tick labels

    # Set font size for category axis (months)
    chart.category_axis.tick_labels.font.size = fontsize  # Set font size for category axis labels
    
    # Find the maximum value across all series
    max_value = 0
    for series in chart.plots[0].series:
        try:
            series_max = max(series.values)
        except:
            series_max = 0
        max_value = max(max_value, series_max)  

    if max_value >= 1000:
        chart.value_axis.tick_labels.number_format = '#,##0'  # add commas to separate thousands
    
    if chart_title:
        chart.chart_title.text_frame.text = title # Set the title text   
        title_font = chart.chart_title.text_frame.paragraphs[0].font
        title_font.bold = True
        title_font.size = fontsize_title
    else:
        chart.has_title = False

    # Remove Gridlines (Line Chart Specific)
    value_axis = chart.value_axis
    category_axis = chart.category_axis
    value_axis.has_major_gridlines = False
    value_axis.has_minor_gridlines = False
    category_axis.has_major_gridlines = False
    category_axis.has_minor_gridlines = False
    
    # if percentage
    if percentage:
        category_axis.tick_labels.NumberFormat = '0"%"'

    return chart
import seaborn as sns
import matplotlib.pyplot as plt
import os

def generate_correlation_plot(df, x_cols, y_cols, brand, category, save_dir="result"):
    if not os.path.exists(save_dir):
        os.makedirs(save_dir)

    correlation_matrix = df.corr()
    sub_corr = correlation_matrix.loc[y_cols, x_cols]

    plt.figure(figsize=(8, 6))
    sns.heatmap(sub_corr, fmt=".2f", annot=True, cmap='RdYlGn', vmin=-1, vmax=1)
    plt.title(f'{brand.upper()} {category.upper()} TV Correlation')
    plt.tight_layout()

    file_path = os.path.join(save_dir, f"{category.lower()}_{brand.lower()}_tbca_correlation.png")
    plt.savefig(file_path)
    plt.close()
    print(f"✅ Saved correlation plot to {file_path}")
    return file_path
def add_correlation_slide(prs, image_path, title_text="📈 Correlation Heatmap", insight_text=None):
   
    slide = prs

    # Add title
    format_title(
        slide,
        text=title_text,
        alignment=PP_ALIGN.LEFT,
        font_name="Calibri",
        font_size=24,
        font_bold=True
    )

    # Add heatmap image
    slide.shapes.add_picture(image_path, left=Inches(1), top=Inches(1.5), width=Inches(6))



    return slide

def column_clustered_chart(slide, df, x, y, cx, cy, legend=True, data_show=False):
    from pptx.chart.data import CategoryChartData
    df=df.T.copy()
    chart_data = CategoryChartData()
    chart_data.categories = list(df.columns)  # Categories = months (or time axis)

    for idx, row in df.iterrows():
        chart_data.add_series(str(idx), row.values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Legend
    chart.has_legend = legend
    if legend:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.font.size = Pt(9)

    # Axis formatting
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)

    # Optional: Show data labels
    if data_show:
        for series in chart.series:
            series.data_labels.show_value = True
            series.data_labels.font.size = Pt(9)
            series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    # Clean gridlines
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False

    return chart


def pie_chart(slide,df,x,y,cx,cy,fontsize=15,title=False,legend_right = True):    
    # Convert the transposed DataFrame into chart data
    chart_data = CategoryChartData()
    # Add the brand names as categories to the chart data
    for i in df.transpose().columns:
        chart_data.add_category(i)

    # Add the SOV values as series to the chart data
    for index, row in df.transpose().iterrows():
        chart_data.add_series(index, row.values)

    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

    if title:
        chart.has_title = True
        #chart.chart_title.text_frame.text = str(chart.chart_title.text_frame.text).upper()
        # chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)  # Set font size to 24pt
        # chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)  # Set font color to black
    else:
        chart.has_title = False
        
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.RIGHT if legend_right else XL_LEGEND_POSITION.TOP 
    # Set legend font size to 10 points
    chart.legend.font.size = Pt(fontsize) 


    for series in chart.plots[0].series:
        for i, val in enumerate(series.values):
            if val == 0:
                series.points[i].data_label.has_text_frame = True
            series.data_labels.show_value = True
            series.data_labels.font.size = Pt(fontsize)
            series.data_labels.number_format = '0%'
            series.data_labels.position = XL_LABEL_POSITION.BEST_FIT
            
    # chart.chart_title.text_frame.text = chart.chart_title.text_frame.text + 'budget distribution'
    # chart.chart_title.text_frame.paragraphs[0].font.size = Pt(fontsize)
    # chart.chart_title.text_frame.color.rgb = RGBColor(0,0,0)
          
    return chart

def table_default(slide, df, left, top, width, height, width_row, height_row, 
                  header=True, upper=False, fontsize=10, alignment=PP_ALIGN.LEFT, percentage=False):
    from pptx.util import Inches, Pt

    # Build table data
    table_data = df.copy()
    if header:
        table_data.columns = [col.upper() if upper else col for col in df.columns]
        table_data = [table_data.columns.tolist()] + table_data.values.tolist()
    else:
        table_data = table_data.values.tolist()

    # Create table
    rows, cols = len(table_data), len(table_data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Apply cell styles
    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            if isinstance(val, (int, float)):
                text = f"{val:.0%}" if percentage else f"{val:,}"
            else:
                text = str(val).upper() if upper else str(val)
            cell.text = text
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.margin_left = 0
            cell.text_frame.margin_right = 0
            cell.text_frame.margin_top = 0
            cell.text_frame.margin_bottom = 0

            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = alignment
                for run in paragraph.runs:
                    run.font.size = Pt(fontsize)

    # Set column widths
    for col_idx, col in enumerate(table.columns):
        if col_idx < len(width_row):
            col.width = width_row[col_idx]

    # Set row heights
    for row in table.rows:
        row.height = height_row

    return table
