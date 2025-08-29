# tbca_ppt_generator.py
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from utils import *
import os
from pptx.chart.data import BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Pt
from pptx.dml.color import RGBColor
def load_data():
    df = pd.read_csv('data/mock_brandbook_monthly.csv')
    df['monthyear'] = pd.to_datetime(df['monthyear'])
    df.set_index('monthyear', inplace=True)
    return df

def generate_ppt():
    # Load template and data
    template_path = 'template/TBCA Analysis Template_blank.pptx'
    output_path = 'output/TBCA_oral_care_pepsodent_core_2024.pptx'
    df = load_data()

    prs = Presentation(template_path)
    slide = prs.slides[5] 

    # Add formatted title
    format_title(
        slide,
        text="Brand Performance Overview",
        alignment=PP_ALIGN.CENTER,
        font_name="Calibri",
        font_size=24,
        font_bold=True,
        font_color=RGBColor(0, 51, 102),
        top=Pt(20),
        width=Pt(900)
    )

    # Generate chart (line marker chart for trend)
    line_marker_chart(
        slide,
        df.transpose(),
        x=Inches(1),
        y=Inches(1.5),
        cx=Inches(8),
        cy=Inches(4.5),
        legend=True,
        legend_position=XL_LEGEND_POSITION.BOTTOM,
        data_show=True,
        chart_title=True,
        title="Monthly Share of Voice",
        percentage=False
    )

    # assuming utils.py functions are imported: format_title, line_marker_chart, pie_chart, column_stacked_chart, etc.

    # Load your presentation template and data first
    prs = Presentation(template_path)
    slide_layout = prs.slide_layouts[5]  # blank slide layout

    # --- Slide 1: Line Marker Chart ---
    slide1 = prs.slides[1]
    format_title(slide1, "Monthly Share of Voice", PP_ALIGN.CENTER, "Calibri", 24, font_bold=True, font_color=RGBColor(0, 51, 102), top=Pt(20), width=Pt(900))
    line_marker_chart(slide1, df.transpose(), Inches(1), Inches(1.5), Inches(8), Inches(4.5), legend=True, legend_position=XL_LEGEND_POSITION.BOTTOM, data_show=True, chart_title=True, title="Monthly Share of Voice")

    # --- Slide 2: Column Stacked Chart ---
    # slide2 = prs.slides[3]
    # format_title(slide2, "Stacked Sales Overview", PP_ALIGN.CENTER, "Calibri", 24, font_bold=True, font_color=RGBColor(0, 51, 102), top=Pt(20), width=Pt(900))
    # column_stacked_chart(slide2, df, Inches(1), Inches(1.5), Inches(8), Inches(4.5))

    # --- Slide 3: Pie Chart ---
    slide3 =prs.slides[2]
    format_title(slide3, "Market Share Distribution", PP_ALIGN.CENTER, "Calibri", 24, font_bold=True, font_color=RGBColor(0, 51, 102), top=Pt(20), width=Pt(900))
    pie_chart(slide3, df, Inches(2), Inches(1.5), Inches(6), Inches(4), fontsize=12, title=True)

    # # --- Slide 4: Table ---
    # slide4 = prs.slides[6]
    # format_title(slide4, "Summary Table", PP_ALIGN.CENTER, "Calibri", 24, font_bold=True, font_color=RGBColor(0, 51, 102), top=Pt(20), width=Pt(900))
    # table_default(slide4, df.head(5), Inches(1), Inches(1.5), Inches(8), Inches(2), width_row=[Inches(2), Inches(2), Inches(2)], height_row=Pt(20), header=True, fontsize=10)

    df_pct = df.div(df.sum(axis=1), axis=0)


    # --- Slide 2: Column Stacked Chart ---
    slide2 = prs.slides[3]
    format_title(
        slide2, "SOV Distribution (Stacked View)",
        alignment=PP_ALIGN.CENTER, font_name="Calibri", font_size=24, font_bold=True
    )
    column_stacked_chart(
        slide2, df_pct.transpose(),
        x=Inches(1), y=Inches(1.5), cx=Inches(8), cy=Inches(4.5)
    )

    # --- Slide 3: Table Slide ---
    slide3 = prs.slides[4]
    format_title(slide3, "🧾 Data Table View", alignment=PP_ALIGN.LEFT, font_name="Calibri", font_size=20)

    table_default(
        slide=slide3,
        df=df.reset_index().tail(5),  # Show last 5 months
        left=Inches(0.5), top=Inches(1.3),
        width=Inches(9), height=Inches(4),
        width_row=[Inches(1.5), Inches(2), Inches(2), Inches(2), Inches(2)],
        height_row=Pt(20),
        header=True, fontsize=10, upper=True
    )


    # --- Clustered Column Chart Slide ---
    slide4 = prs.slides[5]
    format_title(
        slide4,
        "📊 Brand Comparison by Month (Clustered Bar)",
        alignment=PP_ALIGN.CENTER,
        font_name="Calibri",
        font_size=22,
        font_bold=True,
        font_color=RGBColor(0, 51, 102)
    )

    # Transpose so brands = series, months = X-axis
    column_clustered_chart(
        slide=slide4,
        df=df_pct.transpose(),  # brands as rows, months as columns
        x=Inches(1), y=Inches(1.5),
        cx=Inches(8), cy=Inches(4.5),
        legend=True,
        data_show=True
    )
    import pandas as pd
    import json

    json_data = """
    {
    "monthyear": ["2024-01-01", "2024-02-01", "2024-03-01", "2024-04-01", "2024-05-01", "2024-06-01"],
    "TV Spend": [100000, 120000, 110000, 130000, 125000, 115000],
    "Digital Spend": [50000, 60000, 55000, 65000, 62000, 59000],
    "Total GRP": [1200, 1300, 1250, 1400, 1350, 1280],
    "SOV": [0.25, 0.27, 0.26, 0.28, 0.27, 0.26],
    "Top of Mind": [0.60, 0.62, 0.61, 0.64, 0.63, 0.62]
    }
    """
    x_cols = ['TV Spend', 'Digital Spend', 'Total GRP']
    y_cols = ['SOV', 'Top of Mind']
    brand = "Pepsodent"
    category = "Oral Care"
    data_dict = json.loads(json_data)
    df = pd.DataFrame(data_dict)
    df['monthyear'] = pd.to_datetime(df['monthyear'])

    print(df)

    # Generate heatmap PNG
    correlation_image = generate_correlation_plot(df, x_cols, y_cols, brand, category)

    # Create slide with image + AI insight
    add_correlation_slide(
        prs=prs.slides[6],
        image_path=correlation_image,
        title_text=f"{brand} - Media vs Brand Metric Correlation",
    )
    # --- Slide 6: Bar Stacked Chart ---
    slide6 = prs.slides[7]

    format_title(
        slide6,
        text="📊 Media Channel Contribution per Brand (Stacked Bar)",
        alignment=PP_ALIGN.LEFT,
        font_name="Calibri",
        font_size=22,
        font_bold=True,
        font_color=RGBColor(0, 51, 102),
    )

    # Sample DataFrame (replace this with your actual data loading if needed)
    bar_df = pd.DataFrame({
        "ATL": [0.2, 0.25, 0.4],
        "BTL": [0.3, 0.25, 0.2],
        "Digital": [0.5, 0.5, 0.4]
    }, index=["Brand A", "Brand B", "Brand C"])

    # Use the function
    bar_stacked_chart(
        slide=slide6,
        df=bar_df,
        x=Inches(1),
        y=Inches(1.5),
        cx=Inches(8),
        cy=Inches(4),
        percentage=True,
        data_show=True,
        legend=True,
        axis=True,
        fontsize=Pt(10),
        fontsize_label=Pt(9)
    )

    # Dummy data
    data = {
        'Brand A': [10, 20, 15],
        'Brand B': [12, 18, 22],
        'Brand C': [8, 25, 12]
    }
    index = ['Q1', 'Q2', 'Q3']
    df = pd.DataFrame(data, index=index)

    slide = prs.slides[8]# Add formatted title textbox above
    format_title(slide, "Monthly Brand Spend Comparison", PP_ALIGN.CENTER, 'Calibri', 18, True, RGBColor(0,0,0), left=Pt(75), top=Pt(40), width=Pt(850), height=Pt(50))

    # Add the chart to slide
    bar_clustered_chart(
        slide=slide,
        df=df,
        x=Inches(1),
        y=Inches(1.5),
        cx=Inches(7),
        cy=Inches(4),
        legend=True,
        data_show=True,
        value_col='Sales'
    )
    # Save the presentation
    import pandas as pd

    df = pd.DataFrame({
        'Demographics': ['Group A', 'Group B', 'Group C'],
        'x': [10, 20, 30],
        'y': [40, 25, 35],
        'Size': [15, 30, 45],
        'colors': [(255, 0, 0), (0, 255, 0), (0, 0, 255)]  # RGB tuples
    })
    slide = prs.slides[9]
    format_title(slide, "Customer Demographics Bubble Chart", PP_ALIGN.CENTER, 'Calibri', 18, True, RGBColor(0,0,0), left=Pt(75), top=Pt(40), width=Pt(850), height=Pt(50))
    bubble_chart(
        slide, df,
        x=Inches(1), y=Inches(2), cx=Inches(6), cy=Inches(4),
        x_col='x', y_col='y',
        legend=True, data_show=True
    )
    import pandas as pd
    slide = prs.slides[10]
    data = {
        'name': ['Group A', 'Group B', 'Group C', 'Group D', 'Group E'],
        'XValue': [10, 20, 30, 40, 50],
        'YValue': [5, 15, 25, 35, 45]
    }

    df = pd.DataFrame(data)
    print(df)

    # Add a separate title textbox (optional)
    format_title(
        slide, "Customer Demographics Scatter Plot",
        alignment=PP_ALIGN.CENTER,
        font_name="Calibri",
        font_size=18,
        font_bold=True,
        font_color=RGBColor(0, 0, 0),
        left=Pt(75),
        top=Pt(40),
        width=Pt(850),
        height=Pt(50)
    )
    # Add the scatter plot without any chart title
    scatter_plot(
        slide, df,
        x_col='XValue',
        y_col='YValue',
        label_col='name',  # <-- make sure this matches your dataframe
        x=Inches(1), y=Inches(2),
        cx=Inches(6), cy=Inches(4)
    )


    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    generate_ppt()
